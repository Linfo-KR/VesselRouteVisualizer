"""
PPT 자동 생성을 위한 모듈

이 모듈은 pandas DataFrame 형식의 데이터를 입력받아,
사전 정의된 설정(레이아웃, 스타일)에 따라 파워포인트(.pptx) 파일을 생성합니다.

주요 클래스:
- PptAutomator: PPT 생성의 모든 과정을 캡슐화한 메인 클래스.
"""

import pandas as pd
from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from typing import Dict, Any, List
from pathlib import Path
import json


def load_config(config_path: Path) -> Dict[str, Any]:
    """JSON 설정 파일을 로드합니다."""
    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except FileNotFoundError:
        print(f"오류: 설정 파일을 찾을 수 없습니다 - {config_path}")
        raise
    except json.JSONDecodeError:
        print(f"오류: JSON 설정 파일의 형식이 잘못되었습니다 - {config_path}")
        raise

class PptAutomator:
    """
    PPT 자동 생성을 담당하는 클래스.
    설정 객체를 기반으로 동작합니다.
    """

    def __init__(self, config: Dict[str, Any]):
        """
        PptAutomator를 초기화합니다.
        JSON 설정에서 읽어온 색상 배열을 RGBColor 객체로 변환합니다.

        Args:
            config: PPT 레이아웃 및 스타일 정보가 담긴 설정 딕셔너리.
        """
        self.config = config
        # JSON으로부터 읽어온 색상 배열 [R, G, B]를 RGBColor 객체로 변환
        self.colors = {k: RGBColor(*v) for k, v in config["colors"].items()}

    def generate_presentation(self, data: pd.DataFrame) -> None:
        """
        데이터프레임으로부터 전체 PPT 프레젠테이션을 생성합니다.

        Args:
            data: 슬라이드를 생성할 데이터가 담긴 pandas DataFrame.
        """
        prs = Presentation()
        dims = self.config["slide_dimensions"]
        prs.slide_width = Cm(dims["width"])
        prs.slide_height = Cm(dims["height"])

        for _, row in data.iterrows():
            self._add_service_slide(prs, row)

        output_path = self.config["file_paths"]["output_ppt"]
        prs.save(output_path)
        print(f"성공: 프레젠테이션이 '{output_path}' 파일로 저장되었습니다.")

    def _add_service_slide(self, prs: Presentation, data_row: pd.Series) -> None:
        """하나의 항로 데이터에 대한 슬라이드를 생성하고 프레젠테이션에 추가합니다."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # 1. Region Box 생성
        spec = self.config["elements"]["box_region"]
        text = f"{data_row.get('Rno', '')}_{data_row.get('Region', '')}"
        self._create_textbox(slide, spec, text)

        # 2. Service Name Box 생성
        spec = self.config["elements"]["box_service_name"]
        text = f"[{data_row.get('Carriers', '')}] {data_row.get('ServiceName', '')}"
        self._create_textbox(slide, spec, text)
        
        # 3. 서비스 정보 테이블 생성
        self._create_service_info_table(slide, data_row)
        
        # 4. Proforma 테이블들 생성
        self._create_proforma_tables(slide, data_row)

    def _create_textbox(self, slide: Slide, spec: Dict[str, Any], text: str) -> None:
        """설정에 따라 텍스트박스를 생성합니다."""
        shape = slide.shapes.add_textbox(
            Cm(spec["pos"]["left"]), Cm(spec["pos"]["top"]),
            Cm(spec["size"]["width"]), Cm(spec["size"]["height"])
        )
        shape.text = text
        
        # 배경색 채우기
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors[spec["fill"]]
        
        # 텍스트 프레임 및 폰트 설정
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        
        font_spec = spec["font"]
        p.font.name = font_spec["name"]
        p.font.size = Pt(font_spec["size"])
        p.font.bold = font_spec["bold"]
        p.font.color.rgb = self.colors[font_spec["color"]]
        
        if spec["align"] == "CENTER":
            p.alignment = PP_ALIGN.CENTER
        elif spec["align"] == "LEFT":
            p.alignment = PP_ALIGN.LEFT

    def _set_cell_style(self, cell, text: str, font_spec: Dict, font_color_key: str,
                        bg_color_key: str = None, align: str = "CENTER", valign: str = "MIDDLE") -> None:
        """테이블 셀의 스타일을 설정합니다."""
        cell.text = str(text)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE if valign == "MIDDLE" else MSO_VERTICAL_ANCHOR.TOP

        if bg_color_key:
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors[bg_color_key]

        p = cell.text_frame.paragraphs[0]
        p.font.name = font_spec["name"]
        p.font.size = Pt(font_spec["size"])
        p.font.bold = font_spec.get("bold", False)
        p.font.color.rgb = self.colors[font_color_key]
        
        p.alignment = PP_ALIGN.LEFT if align == "LEFT" else PP_ALIGN.CENTER

    def _create_service_info_table(self, slide: Slide, data_row: pd.Series) -> None:
        """서비스 정보 테이블을 생성합니다."""
        spec = self.config["elements"]["tbl_service_info"]
        pos, size = spec["pos"], spec["size"]
        
        cols_spec = spec["columns"]
        col_keys = list(cols_spec.keys())
        
        table = slide.shapes.add_table(2, len(col_keys), Cm(pos["left"]), Cm(pos["top"]), Cm(size["width"]), Cm(size["height"])).table

        for i, key in enumerate(col_keys):
            table.columns[i].width = Cm(cols_spec[key]["width"])
            # Header
            self._set_cell_style(table.cell(0, i), key.replace("PortRotation", "Rotation").replace("Dur", "Dur."),
                                 spec["header_font"], "WHITE", "TABLE_HEADER_BG")
            # Body
            self._set_cell_style(table.cell(1, i), data_row.get(key, ""),
                                 spec["body_font"], "BLACK", align=cols_spec[key]["align"])

    def _create_proforma_tables(self, slide: Slide, data_row: pd.Series) -> None:
        """Proforma 정보 테이블들을 생성합니다."""
        spec = self.config["elements"]["tbl_proforma"]
        pos, size = spec["pos"], spec["size"]
        
        for i in range(1, spec["max_count"] + 1):
            t_name = data_row.get(f"T{i}Name")
            if pd.isna(t_name) or str(t_name).strip() == "":
                continue

            left = pos["left"] + (size["width"] + spec["gap"]) * (i - 1)
            table = slide.shapes.add_table(len(spec["rows"]), 2, Cm(left), Cm(pos["top"]), Cm(size["width"]), Cm(size["height"])).table
            
            table.columns[0].width = Cm(spec["columns"]["header"]["width"])
            table.columns[1].width = Cm(spec["columns"]["body"]["width"])

            for row_idx, header_text in enumerate(spec["rows"]):
                # Header
                self._set_cell_style(table.cell(row_idx, 0), header_text,
                                     spec["header_font"], "WHITE", "PROFORMA_HEADER_BG")
                # Body
                data_key = {"Terminal": f"T{i}Name", "Weekly throughput": f"T{i}Wtp", "Schedule": f"T{i}Sch"}[header_text]
                self._set_cell_style(table.cell(row_idx, 1), data_row.get(data_key, ""),
                                     spec["body_font"], "BLACK")
