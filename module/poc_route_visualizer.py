"""
항로 시각화 PoC(Proof of Concept) 모듈
"""
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
import json
from pathlib import Path

class RouteVisualizer:
    """
    CSV 데이터와 좌표계를 기반으로 PPT에 항로를 시각화하는 클래스.
    """

    def __init__(self, port_coords_path: Path):
        """
        RouteVisualizer를 초기화합니다.

        Args:
            port_coords_path: 항구 좌표 정보가 담긴 JSON 파일 경로.
        """
        try:
            with open(port_coords_path, 'r', encoding='utf-8') as f:
                self.port_coords = json.load(f)["ports"]
        except FileNotFoundError:
            print(f"오류: 좌표 파일을 찾을 수 없습니다 - {port_coords_path}")
            raise
        self.port_style = {
            "fill_color": RGBColor(0, 112, 192),
            "line_color": RGBColor(255, 255, 255),
            "radius": 0.25
        }
        self.text_style = {
            "font_size": Pt(8),
            "font_name": "Arial"
        }
        self.connector_style = {
            "color": RGBColor(0, 32, 96),
            "width": Pt(1.5)
        }

    def draw_route(self, slide, port_rotation_str: str):
        """
        하나의 항로(Port Rotation)를 슬라이드에 그립니다.

        Args:
            slide: 도형을 그릴 pptx.slide.Slide 객체.
            port_rotation_str: 쉼표로 구분된 항구 이름 문자열.
        """
        if not port_rotation_str or pd.isna(port_rotation_str):
            return

        port_names = [p.strip() for p in port_rotation_str.split(',')]
        
        route_points = []
        for name in port_names:
            if name in self.port_coords:
                route_points.append({
                    "name": name,
                    "coords": self.port_coords[name]
                })
            else:
                print(f"경고: '{name}' 항구의 좌표를 찾을 수 없습니다. 건너뜁니다.")

        # 항구 아이콘 및 이름 그리기
        for point in route_points:
            self._draw_port(slide, point["name"], point["coords"]["x"], point["coords"]["y"])

        # 항로 연결선 그리기
        for i in range(len(route_points) - 1):
            p1 = route_points[i]
            p2 = route_points[i+1]
            self._draw_connector(slide, p1["coords"]["x"], p1["coords"]["y"], p2["coords"]["x"], p2["coords"]["y"])

    def _draw_port(self, slide, name: str, x: float, y: float):
        """슬라이드에 항구 아이콘과 이름을 그립니다."""
        radius = self.port_style["radius"]
        # 원의 중심이 (x, y)가 되도록 좌상단 좌표 조정
        left = Cm(x - radius)
        top = Cm(y - radius)
        dia = Cm(radius * 2)

        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, dia, dia)
        
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.port_style["fill_color"]
        shape.line.color.rgb = self.port_style["line_color"]
        shape.line.width = Pt(1)

        # 텍스트 박스 (항구 이름)
        # 아이콘 아래에 위치하도록 top 위치 조정
        text_box = slide.shapes.add_textbox(left, Cm(y + radius * 1.5), Cm(2.5), Cm(0.5))
        p = text_box.text_frame.paragraphs[0]
        p.text = name
        p.font.name = self.text_style["font_name"]
        p.font.size = self.text_style["font_size"]

    def _draw_connector(self, slide, x1: float, y1: float, x2: float, y2: float):
        """두 지점 사이에 화살표가 있는 연결선을 그립니다."""
        # add_connector는 시작점과 끝점의 좌표를 직접 사용하지 않고,
        # 시작 도형과 끝 도형을 연결하는 방식이 일반적이나, 여기서는 좌표로 직접 선을 그림.
        # add_shape으로 직선을 그리고 화살표를 추가하는 방식으로 구현
        line_shape = slide.shapes.add_connector(
            1, # 1 for MSO_CONNECTOR.STRAIGHT
            Cm(x1), Cm(y1),
            Cm(x2), Cm(y2)
        )
        line_shape.line.color.rgb = self.connector_style["color"]
        line_shape.line.width = self.connector_style["width"]
        
        # 화살표 설정
        line_shape.line.end_arrowhead_style = 2 # 2 for MSO_ARROWHEAD.TRIANGLE

    def generate_poc_presentation(self, csv_path: Path, template_path: Path, output_path: Path, num_rows: int = 5):
        """
        CSV 파일에서 데이터를 읽고, 템플릿을 기반으로 PoC용 프레젠테이션을 생성합니다.

        Args:
            csv_path: 입력 데이터 CSV 파일 경로.
            template_path: 사용할 PPT 템플릿 파일 경로.
            output_path: 저장할 PPT 파일 경로.
            num_rows: 처리할 데이터 행의 수.
        """
        try:
            df = pd.read_csv(csv_path, nrows=num_rows)
        except FileNotFoundError:
            print(f"오류: CSV 파일을 찾을 수 없습니다 - {csv_path}")
            return

        try:
            prs = Presentation(template_path)
        except FileNotFoundError:
            print(f"오류: 템플릿 파일을 찾을 수 없습니다 - {template_path}")
            return

        # 템플릿의 첫 번째 슬라이드 레이아웃을 사용
        slide_layout = prs.slide_layouts[0] if prs.slide_layouts else prs.slide_layouts[6]
        if prs.slides:
            slide_layout = prs.slides[0].slide_layout


        for _, row in df.iterrows():
            slide = prs.slides.add_slide(slide_layout)
            
            # 슬라이드 제목 추가 (템플릿에 제목 영역이 있다고 가정하고 업데이트 시도, 없으면 새로 추가)
            try:
                title_shape = slide.shapes.title
                title_shape.text = f"[{row.get('Carriers', '')}] {row.get('ServiceName', '')}"
            except AttributeError:
                # 템플릿에 기본 제목 도형이 없을 경우, 새로 추가
                title_shape = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(30), Cm(1.5))
                title_shape.text = f"[{row.get('Carriers', '')}] {row.get('ServiceName', '')}"
                title_shape.text_frame.paragraphs[0].font.size = Pt(24)

            # 항로 시각화
            self.draw_route(slide, row['PortRotation'])

        prs.save(output_path)
        print(f"성공: PoC 프레젠테이션이 '{output_path}' 파일로 저장되었습니다.")


if __name__ == '__main__':
    # 프로젝트 루트 경로를 기준으로 파일 경로 설정
    ROOT_DIR = Path(__file__).parent.parent
    
    COORDS_PATH = ROOT_DIR / "config" / "port_coordinates.json"
    INPUT_CSV_PATH = ROOT_DIR / "data" / "test" / "proforma_test.csv"
    TEMPLATE_PPT_PATH = ROOT_DIR / "data" / "test" / "pptx_test.pptx"
    OUTPUT_PPT_PATH = ROOT_DIR / "data" / "output" / "route_visualization_poc_template.pptx"
    
    print("템플릿 기반 항로 시각화 PoC를 시작합니다...")
    print(f"좌표 파일: {COORDS_PATH}")
    print(f"입력 CSV: {INPUT_CSV_PATH}")
    print(f"템플릿 PPT: {TEMPLATE_PPT_PATH}")
    print(f"출력 PPT: {OUTPUT_PPT_PATH}")
    
    visualizer = RouteVisualizer(COORDS_PATH)
    visualizer.generate_poc_presentation(INPUT_CSV_PATH, TEMPLATE_PPT_PATH, OUTPUT_PPT_PATH)
